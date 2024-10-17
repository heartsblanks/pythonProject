import os
from pptx import Presentation
import comtypes.client

def convert_ppt_to_pptx(folder_path):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    for filename in os.listdir(folder_path):
        if filename.endswith(".ppt"):
            full_path = os.path.join(folder_path, filename)
            pptx_path = os.path.join(folder_path, filename.replace(".ppt", ".pptx"))

            try:
                # Open and save as .pptx
                presentation = powerpoint.Presentations.Open(full_path)
                presentation.SaveAs(pptx_path, 24)  # 24 is the format ID for .pptx
                presentation.Close()
                
                # Delete the original .ppt file after successful conversion
                os.remove(full_path)
                print(f"Converted {filename} to .pptx and removed the original .ppt file.")
            except Exception as e:
                print(f"Failed to convert {filename}: {e}")

    powerpoint.Quit()

def analyze_slide_flow(folder_path):
    results = []

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            presentation = Presentation(file_path)
            slide_data = {
                "filename": filename,
                "connections": []
            }

            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_shapes = []
                slide_arrows = []

                # Extract shapes and arrows
                for shape in slide.shapes:
                    if shape.is_placeholder or shape.has_text_frame:
                        # Collect shapes with text (e.g., process blocks)
                        slide_shapes.append({
                            "name": shape.name,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                            "text": shape.text_frame.text if shape.has_text_frame else ""
                        })
                    elif shape.shape_type == 3:  # Type 3 indicates a line (potential arrow)
                        # Assume these are arrows
                        slide_arrows.append({
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                            "rotation": shape.rotation
                        })

                # Attempt to connect arrows to shapes
                for arrow in slide_arrows:
                    arrow_start_x = arrow["left"]
                    arrow_start_y = arrow["top"]
                    arrow_end_x = arrow["left"] + arrow["width"]
                    arrow_end_y = arrow["top"] + arrow["height"]

                    for shape in slide_shapes:
                        # Check if arrow is connected to the shape by comparing coordinates
                        shape_left = shape["left"]
                        shape_top = shape["top"]
                        shape_right = shape["left"] + shape["width"]
                        shape_bottom = shape["top"] + shape["height"]

                        # Check if arrow ends near shape boundaries (basic overlap detection)
                        if (shape_left - 100 < arrow_end_x < shape_right + 100 and
                            shape_top - 100 < arrow_end_y < shape_bottom + 100):
                            connection = {
                                "from_shape": shape["text"],
                                "to_shape": "",  # Placeholder for connected shape text
                                "direction": "unknown"  # Placeholder for direction
                            }
                            # Add connection to slide data
                            slide_data["connections"].append(connection)

                results.append(slide_data)

    # Display results
    for slide in results:
        print(f"File: {slide['filename']}")
        for conn in slide["connections"]:
            print(f"  Connection: From '{conn['from_shape']}' to '{conn['to_shape']}' with direction '{conn['direction']}'")

# Usage example
folder_path = "/path/to/your/ppt/folder"  # Replace with your folder path
convert_ppt_to_pptx(folder_path)  # First, convert .ppt to .pptx if needed
analyze_slide_flow(folder_path)  # Then, analyze slide flow direction